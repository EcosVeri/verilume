"""Local document ingestion into Chroma."""

from __future__ import annotations

import hashlib
import io
import json
import logging
import os
import re
import shutil
import tempfile
import unicodedata
from abc import ABC, abstractmethod
from dataclasses import asdict
from contextlib import contextmanager
from collections.abc import Callable
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor, as_completed
from functools import lru_cache
from pathlib import Path

try:
    import fcntl
except ImportError:  # pragma: no cover - non-POSIX fallback
    fcntl = None

from verilume.core.embeddings import EmbeddingService
from verilume.core.equation_repair import repair_math_text
from verilume.core.formula_extraction import extract_formulas, formula_to_text
from verilume.core.formula_store import FormulaStore
from verilume.core.ocr_blocks import OCRBlock, OCRBlockStore, page_text_block
from verilume.core.ocr_cleaning import clean_ocr_text
from verilume.core.retrieval import ChromaRetriever
from verilume.core.schemas import DocumentChunk, DocumentMetadata, IngestResult
from verilume.core.structured_document_store import StructuredDocumentStore
from verilume.core.structured_ocr import StructuredDocument, extract_structured_document
from verilume.settings import AppSettings, ensure_app_dirs

LOGGER = logging.getLogger(__name__)

_PDF_PHONE_ICON_FRAGMENT_RE = re.compile(r"(?:(?<=^)|(?<=\|))\s*/[A-Za-z]{1,6}(?=\+\d)")
_PDF_DELIMITED_ICON_FRAGMENT_RE = re.compile(
    r"(?:(?<=^)|(?<=\|))\s*/(?!/)[^\s|]{1,12}(?=(?:\s*\||\s*$))"
)
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp"}
_PRESENTATION_EXTENSIONS = {".pptx", ".pptm", ".ppsx", ".potx"}
_OCR_CONFIDENCE_FLOOR = 0.35
_PDF_OCR_RENDER_SCALE = 1.0

ProgressCallback = Callable[[str, int, int], None]


class FileHandler(ABC):
    """Extract text from one or more file types."""

    extensions: set[str] = set()

    @abstractmethod
    def extract_pages(self, path: Path) -> tuple[list[tuple[int | None, str]], int]:
        """Return extracted text by page and PDF page count when available."""


class PdfHandler(FileHandler):
    extensions = {".pdf"}

    def extract_pages(self, path: Path) -> tuple[list[tuple[int | None, str]], int]:
        return _extract_pdf(path)


class DocxHandler(FileHandler):
    extensions = {".docx"}

    def extract_pages(self, path: Path) -> tuple[list[tuple[int | None, str]], int]:
        return _extract_docx(path)


class TextHandler(FileHandler):
    extensions = {".txt", ".md", ".markdown", ".csv"}

    def extract_pages(self, path: Path) -> tuple[list[tuple[int | None, str]], int]:
        return [(None, _read_text_file(path))], 0


class PptxHandler(FileHandler):
    extensions = _PRESENTATION_EXTENSIONS

    def extract_pages(self, path: Path) -> tuple[list[tuple[int | None, str]], int]:
        return _extract_pptx(path)


class ImageHandler(FileHandler):
    extensions = _IMAGE_EXTENSIONS

    def extract_pages(self, path: Path) -> tuple[list[tuple[int | None, str]], int]:
        return _extract_image(path)


FILE_HANDLERS: dict[str, FileHandler] = {}
SUPPORTED_EXTENSIONS: set[str] = set()


def register_file_handler(handler: FileHandler) -> None:
    for extension in handler.extensions:
        key = extension.lower()
        FILE_HANDLERS[key] = handler
        SUPPORTED_EXTENSIONS.add(key)


def supported_extensions() -> set[str]:
    return set(FILE_HANDLERS)


register_file_handler(PdfHandler())
register_file_handler(DocxHandler())
register_file_handler(TextHandler())
register_file_handler(PptxHandler())
register_file_handler(ImageHandler())


class ReadonlyChromaError(RuntimeError):
    """Raised when Chroma's SQLite store rejects writes."""


class IngestStateError(RuntimeError):
    """Raised when ingestion would leave the knowledge base degraded."""


def save_uploaded_file(file_name: str, data: bytes, docs_dir: Path) -> Path:
    docs_dir.mkdir(parents=True, exist_ok=True)
    clean_name = Path(file_name).name
    target = docs_dir / clean_name
    target.write_bytes(data)
    return target


def removable_documents(docs_dir: Path) -> list[str]:
    if not docs_dir.exists():
        return []
    return sorted(str(path.relative_to(docs_dir)) for path in supported_files(docs_dir))


def remove_documents(settings: AppSettings, documents: list[str]) -> list[str]:
    ensure_app_dirs(settings)
    if not documents:
        return []

    with _ingest_lock(settings):
        docs_dir = settings.docs_dir
        docs_root = docs_dir.resolve()
        manifest = load_manifest(settings.manifest_path)
        ingestor = DocumentIngestor(settings)
        removed: list[str] = []

        try:
            for document in documents:
                relative = Path(document)
                raw_path = docs_dir / relative
                path = raw_path.resolve()
                if docs_root not in path.parents and path != docs_root:
                    continue

                manifest.pop(str(raw_path), None)
                manifest.pop(str(path), None)
                ingestor.retriever.delete_document(str(raw_path))
                specialized_delete = getattr(ingestor, "_delete_specialized_document", None)
                if callable(specialized_delete):
                    specialized_delete(path.name)
                if path != raw_path:
                    ingestor.retriever.delete_document(str(path))
                try:
                    if path.exists():
                        path.unlink()
                    removed.append(str(relative))
                except OSError:
                    continue

            write_manifest(settings.manifest_path, manifest)
            return removed
        finally:
            ingestor.retriever.close(clear_system_cache=True)


def supported_files(docs_dir: Path) -> list[Path]:
    if not docs_dir.exists():
        return []
    return sorted(
        path
        for path in docs_dir.rglob("*")
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
    )


def file_hash(path: Path) -> str:
    digest = hashlib.blake2b(digest_size=20)
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def load_manifest(path: Path) -> dict[str, dict]:
    if not path.exists():
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}


def write_manifest(path: Path, manifest: dict[str, dict]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(manifest, indent=2, sort_keys=True), encoding="utf-8")


def document_metadata_from_manifest(settings: AppSettings) -> list[DocumentMetadata]:
    manifest = load_manifest(settings.manifest_path)
    documents: list[DocumentMetadata] = []
    for source_path, entry in manifest.items():
        if not isinstance(entry, dict):
            continue
        raw_metadata = entry.get("document_metadata")
        if not isinstance(raw_metadata, dict):
            raw_metadata = _legacy_manifest_document_metadata(source_path, entry)
        metadata = _document_metadata_from_dict(source_path, raw_metadata, entry)
        if metadata:
            documents.append(metadata)
    return documents


def _read_text_file(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


@lru_cache(maxsize=1)
def _ocr_engine():
    from rapidocr_onnxruntime import RapidOCR

    return RapidOCR()


def _normalize_raw_text(text: str) -> str:
    normalized = unicodedata.normalize("NFKC", text or "")
    normalized = normalized.replace("\u00ad", "")
    normalized = normalized.replace("\r\n", "\n").replace("\r", "\n")
    normalized = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", normalized)
    return normalized


def _normalize_extracted_text(text: str) -> str:
    normalized = _normalize_raw_text(text)
    cleaned_lines: list[str] = []
    for raw_line in normalized.splitlines():
        line = re.sub(r"\s+", " ", raw_line).strip()
        if not line:
            if cleaned_lines and cleaned_lines[-1] != "":
                cleaned_lines.append("")
            continue
        cleaned_lines.append(line)
    return "\n".join(cleaned_lines).strip()


def _extract_pdf(path: Path) -> tuple[list[tuple[int | None, str]], int]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[tuple[int | None, str]] = []
    pdfium_document = None
    for index, page in enumerate(reader.pages, start=1):
        text = _normalize_pdf_text(page.extract_text() or "")
        if _pdf_page_needs_ocr(text):
            if pdfium_document is None:
                import pypdfium2 as pdfium

                pdfium_document = pdfium.PdfDocument(str(path))
            ocr_text = _ocr_pdf_page(pdfium_document, index - 1)
            if ocr_text:
                text = ocr_text
        if text.strip():
            pages.append((index, text))
    if pdfium_document is not None:
        try:
            pdfium_document.close()
        except Exception:
            pass
    return pages, len(reader.pages)


def _normalize_pdf_text(text: str) -> str:
    normalized = _normalize_raw_text(text)
    normalized = re.sub(r"(\w)-\n(?=\w)", r"\1", normalized)

    cleaned_lines: list[str] = []
    for raw_line in normalized.splitlines():
        line = raw_line.strip()
        if not line:
            if cleaned_lines and cleaned_lines[-1] != "":
                cleaned_lines.append("")
            continue
        line = _PDF_PHONE_ICON_FRAGMENT_RE.sub(" ", line)
        line = _PDF_DELIMITED_ICON_FRAGMENT_RE.sub(" ", line)
        line = re.sub(r"\s*\|\s*", " | ", line)
        line = re.sub(r"(?:\s\|\s){2,}", " | ", line)
        line = re.sub(r"\s{2,}", " ", line).strip(" |")
        if line:
            cleaned_lines.append(line)

    return "\n".join(cleaned_lines).strip()


def _pdf_page_needs_ocr(text: str) -> bool:
    normalized = re.sub(r"\s+", " ", text or "").strip()
    if not normalized:
        return True
    alnum_count = len(re.findall(r"[A-Za-z0-9]", normalized))
    word_count = len(re.findall(r"\b[A-Za-z0-9][A-Za-z0-9'_-]*\b", normalized))
    return alnum_count <= 2 and word_count == 0


def _ocr_pdf_page(pdfium_document, page_index: int) -> str:
    rendered = pdfium_document[page_index].render(scale=_PDF_OCR_RENDER_SCALE).to_pil()
    try:
        return _ocr_pil_image(rendered)
    finally:
        try:
            rendered.close()
        except Exception:
            pass


def _extract_docx(path: Path) -> tuple[list[tuple[int | None, str]], int]:
    from docx import Document

    document = Document(str(path))
    parts: list[str] = []
    parts.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                parts.append(" | ".join(values))
    return [(None, "\n".join(parts))], 0


def _extract_pptx(path: Path) -> tuple[list[tuple[int | None, str]], int]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    pages: list[tuple[int | None, str]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        _collect_presentation_shapes_text(slide.shapes, parts)
        notes = _presentation_notes_text(slide)
        if notes:
            parts.append(notes)
        text = _normalize_extracted_text("\n".join(parts))
        if text:
            pages.append((slide_number, text))
    return pages, 0


def _collect_presentation_shapes_text(shapes, parts: list[str]) -> None:
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            text = _normalize_extracted_text(getattr(shape, "text", ""))
            if text:
                parts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if values:
                    parts.append(" | ".join(values))
        if hasattr(shape, "image"):
            ocr_text = _ocr_image_blob(shape.image.blob)
            if ocr_text:
                parts.append(ocr_text)
        if hasattr(shape, "shapes"):
            _collect_presentation_shapes_text(shape.shapes, parts)


def _presentation_notes_text(slide) -> str:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""
    parts: list[str] = []
    for shape in notes_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _normalize_extracted_text(getattr(shape, "text", ""))
        if not text or text.lower() == "click to add notes":
            continue
        parts.append(text)
    return _normalize_extracted_text("\n".join(parts))


def _extract_image(path: Path) -> tuple[list[tuple[int | None, str]], int]:
    from PIL import Image, ImageSequence

    pages: list[tuple[int | None, str]] = []
    with Image.open(str(path)) as image:
        frames = [frame.copy() for frame in ImageSequence.Iterator(image)] if getattr(image, "is_animated", False) else [image.copy()]
    total_frames = len(frames)
    for index, frame in enumerate(frames, start=1):
        try:
            text = _ocr_pil_image(frame)
        finally:
            frame.close()
        if text:
            pages.append(((index if total_frames > 1 else 1), text))
    return pages, 0


def _ocr_image_blob(blob: bytes) -> str:
    from PIL import Image

    with Image.open(io.BytesIO(blob)) as image:
        return _ocr_pil_image(image)


def _ocr_pil_image(image) -> str:
    import numpy as np
    from PIL import Image, ImageOps

    prepared = ImageOps.exif_transpose(image).convert("RGB")
    if max(prepared.size) < 1200:
        scale = max(2, int(1200 / max(1, max(prepared.size))))
        prepared = prepared.resize(
            (prepared.width * scale, prepared.height * scale),
            Image.Resampling.LANCZOS,
        )
    prepared = ImageOps.autocontrast(prepared)
    result, _ = _ocr_engine()(np.asarray(prepared))
    return _ocr_result_text(result)


def _ocr_result_text(result) -> str:
    lines: list[str] = []
    for item in result or []:
        if not isinstance(item, (list, tuple)) or len(item) < 2:
            continue
        text = str(item[1] or "").strip()
        if not text:
            continue
        confidence = float(item[2]) if len(item) > 2 else 1.0
        if confidence >= _OCR_CONFIDENCE_FLOOR:
            lines.append(text)
    return _normalize_extracted_text("\n".join(lines))


def extract_pages(path: Path) -> tuple[list[tuple[int | None, str]], int]:
    handler = FILE_HANDLERS.get(path.suffix.lower())
    if handler is None:
        raise ValueError(f"Unsupported file type: {path.suffix}")
    return handler.extract_pages(path)


def chunk_text(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    return chunk_text_semantic(text, chunk_size, chunk_overlap)


def chunk_text_semantic(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    text = "\n".join(line.rstrip() for line in (text or "").splitlines()).strip()
    if not text:
        return []

    paragraphs = [part.strip() for part in text.split("\n\n") if part.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if len(paragraph) > chunk_size:
            if current:
                chunks.append(current.strip())
                current = ""
            chunks.extend(_chunk_by_sentences(paragraph, chunk_size, chunk_overlap))
            continue
        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(candidate) <= chunk_size:
            current = candidate
        else:
            chunks.append(current.strip())
            current = paragraph
    if current:
        chunks.append(current.strip())
    return [chunk for chunk in chunks if chunk.strip()]


def chunk_text_by_paragraphs(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    text = "\n".join(line.rstrip() for line in (text or "").splitlines()).strip()
    if not text:
        return []

    paragraphs = [part.strip() for part in text.split("\n\n") if part.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if len(paragraph) > chunk_size:
            if current:
                chunks.append(current.strip())
                current = ""
            chunks.extend(_sliding_chunks(paragraph, chunk_size, chunk_overlap))
            continue
        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(candidate) <= chunk_size:
            current = candidate
        else:
            chunks.append(current.strip())
            current = paragraph
    if current:
        chunks.append(current.strip())
    return [chunk for chunk in chunks if chunk.strip()]


def _split_sentences(text: str) -> list[str]:
    sentences = re.split(r"(?<=[.!?])\s+(?=[A-Z0-9])", text or "")
    return [sentence.strip() for sentence in sentences if sentence.strip()]


def _chunk_by_sentences(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    sentences = _split_sentences(text)
    if len(sentences) <= 1:
        return _sliding_chunks(text, chunk_size, chunk_overlap)

    chunks: list[str] = []
    current: list[str] = []
    current_length = 0
    for sentence in sentences:
        sentence_length = len(sentence)
        if sentence_length > chunk_size:
            if current:
                chunks.append(" ".join(current).strip())
                current = []
                current_length = 0
            chunks.extend(_sliding_chunks(sentence, chunk_size, chunk_overlap))
            continue
        if current and current_length + sentence_length + 1 > chunk_size:
            chunks.append(" ".join(current).strip())
            current, current_length = _overlap_sentences(current, chunk_overlap)
        current.append(sentence)
        current_length += sentence_length + 1
    if current:
        chunks.append(" ".join(current).strip())
    return [chunk for chunk in chunks if chunk.strip()]


def _overlap_sentences(sentences: list[str], chunk_overlap: int) -> tuple[list[str], int]:
    if chunk_overlap <= 0:
        return [], 0
    overlap: list[str] = []
    length = 0
    for sentence in reversed(sentences):
        next_length = length + len(sentence) + 1
        if next_length > chunk_overlap:
            break
        overlap.insert(0, sentence)
        length = next_length
    return overlap, length


def _sliding_chunks(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    chunks: list[str] = []
    start = 0
    step = max(1, chunk_size - chunk_overlap)
    while start < len(text):
        end = min(len(text), start + chunk_size)
        chunks.append(text[start:end].strip())
        if end == len(text):
            break
        start += step
    return chunks


def _parse_file(
    path: Path, digest: str, settings: AppSettings
) -> tuple[Path, list[DocumentChunk], int, dict, list, list, list]:
    pages, pdf_pages = extract_pages(path)
    document_metadata = _extract_document_metadata(path, pages, settings)
    chunks: list[DocumentChunk] = []
    formula_items = []
    ocr_blocks: list[OCRBlock] = []
    structured_documents: list[StructuredDocument] = []
    for page, text in pages:
        cleaned_text = clean_ocr_text(text)
        if cleaned_text:
            ocr_blocks.append(page_text_block(path.name, page, cleaned_text))
        page_formulas = extract_formulas(
            cleaned_text,
            document=path.name,
            page=page,
            threshold=float(getattr(settings, "formula_detection_threshold", 0.55)),
        )
        formula_items.extend(page_formulas)
        structured_document = extract_structured_document(
            cleaned_text,
            document=path.name,
            page=page,
        )
        if structured_document and structured_document.fields:
            structured_documents.append(structured_document)
        text = repair_math_text(cleaned_text)
        chunker = chunk_text_semantic
        if getattr(settings, "chunk_strategy", "semantic") in {"paragraph", "legacy"}:
            chunker = chunk_text_by_paragraphs
        for chunk in chunker(text, settings.chunk_size, settings.chunk_overlap):
            chunk_index = len(chunks)
            chunk_id = f"{digest}-{page or 0}-{chunk_index}"
            chunk_metadata = {
                "extension": path.suffix.lower(),
                "source_path": str(path),
                **document_metadata,
            }
            section_heading = _section_heading(chunk)
            if section_heading:
                chunk_metadata["section_heading"] = section_heading
            chunks.append(
                DocumentChunk(
                    chunk_id=chunk_id,
                    text=chunk,
                    source_path=path,
                    document=path.name,
                    page=page,
                    chunk_index=chunk_index,
                    file_hash=digest,
                    metadata=chunk_metadata,
                )
            )
        for formula_item in page_formulas:
            chunk_index = len(chunks)
            chunks.append(
                DocumentChunk(
                    chunk_id=f"formula-{formula_item.formula_id}",
                    text=formula_to_text(formula_item),
                    source_path=path,
                    document=path.name,
                    page=page,
                    chunk_index=chunk_index,
                    file_hash=digest,
                    metadata={
                        "extension": path.suffix.lower(),
                        "source_path": str(path),
                        **document_metadata,
                        "content_type": "formula",
                        "formula_id": formula_item.formula_id,
                        "formula_type": formula_item.formula_type or "unknown",
                    },
                )
            )
    manifest_metadata = _build_document_metadata(
        path,
        pages,
        pdf_pages,
        len(chunks),
        document_metadata,
    )
    return path, chunks, pdf_pages, manifest_metadata, formula_items, ocr_blocks, structured_documents


def _extract_document_metadata(
    path: Path,
    pages: list[tuple[int | None, str]],
    settings: AppSettings | None = None,
) -> dict[str, str]:
    sample = "\n".join(text for _page, text in pages[:3])[:12000]
    metadata: dict[str, str] = {}
    title = _infer_document_title(path, sample)
    authors = _infer_authors(sample)
    abstract = _extract_abstract(sample, settings)
    keywords = _extract_keywords(sample, settings)
    document_kind = _infer_document_kind(path, sample)
    for key, value in {
        "document_title": title,
        "authors": authors,
        "abstract": abstract,
        "keywords": keywords,
        "document_kind": document_kind,
    }.items():
        cleaned = _metadata_text(value)
        if cleaned:
            metadata[key] = cleaned
    return metadata


def _infer_document_title(path: Path, text: str) -> str:
    lines = _metadata_lines(text)
    for line in lines[:18]:
        lowered = line.lower()
        if lowered.startswith(("abstract", "keywords", "doi", "http")):
            continue
        if 8 <= len(line) <= 180 and len(line.split()) >= 2:
            return line
    return path.stem.replace("_", " ").replace("-", " ").strip()


def _infer_authors(text: str) -> str:
    lines = _metadata_lines(text)
    candidates: list[str] = []
    for line in lines[1:10]:
        lowered = line.lower()
        if lowered.startswith(("abstract", "keywords", "introduction")):
            break
        if re.search(r"\b(?:university|department|faculty|school|http|doi)\b", lowered):
            continue
        if re.search(r"[A-Z][a-z]+(?:\s+[A-Z][a-z]+)+", line):
            candidates.append(line)
    return "; ".join(candidates[:3])


def _extract_abstract(text: str, settings: AppSettings | None = None) -> str:
    pattern = (
        getattr(settings, "metadata_abstract_pattern", None)
        or AppSettings().metadata_abstract_pattern
    )
    limit = int(
        getattr(settings, "metadata_abstract_limit", None)
        or AppSettings().metadata_abstract_limit
    )
    match = re.search(
        pattern,
        text or "",
    )
    if not match:
        return ""
    return _metadata_text(match.group("abstract"), limit=limit)


def _extract_keywords(text: str, settings: AppSettings | None = None) -> str:
    pattern = (
        getattr(settings, "metadata_keywords_pattern", None)
        or AppSettings().metadata_keywords_pattern
    )
    limit = int(
        getattr(settings, "metadata_keywords_limit", None)
        or AppSettings().metadata_keywords_limit
    )
    match = re.search(pattern, text or "")
    if not match:
        match = re.search(r"(?im)^\s*keywords?\s*[:\-]\s*(?P<keywords>.+)$", text or "")
    return _metadata_text(match.group("keywords"), limit=limit) if match else ""


def _infer_document_kind(path: Path, text: str) -> str:
    haystack = f"{path.name} {text[:4000]}".lower()
    suffix = path.suffix.lower()
    if suffix in _PRESENTATION_EXTENSIONS:
        return "presentation"
    if suffix == ".csv":
        return "spreadsheet"
    if suffix in _IMAGE_EXTENSIONS:
        return "image_document"
    if any(marker in haystack for marker in ("passport", "identity card", "id card", "nationality")):
        return "identity_document"
    if any(marker in haystack for marker in ("invoice", "line item", "total due", "amount due", "vat")):
        return "invoice"
    if any(marker in haystack for marker in ("doctoral thesis", "phd thesis", "dissertation")):
        return "thesis"
    if any(marker in haystack for marker in ("abstract", "journal", "doi", "references", "methodology")):
        return "scientific_paper"
    if any(marker in haystack for marker in ("certificate", "attestation", "exam payment")):
        return "certificate"
    if any(marker in haystack for marker in ("textbook", "worked example", "exercise", "chapter")):
        return "textbook"
    if any(marker in haystack for marker in ("manual", "guide", "handbook", "training", "certification prep")):
        return "manual"
    if any(marker in haystack for marker in ("policy", "procedure", "compliance")):
        return "policy"
    if any(marker in haystack for marker in ("contract", "agreement", "terms and conditions")):
        return "contract"
    if any(marker in haystack for marker in ("meeting minutes", "attendees", "action items")):
        return "meeting_minutes"
    if any(marker in haystack for marker in ("report", "recommendations", "findings")):
        return "report"
    if any(marker in haystack for marker in ("dear ", "sincerely", "letter")):
        return "letter"
    if any(marker in haystack for marker in ("notes", "notebook")):
        return "notes"
    if "book" in haystack:
        return "book"
    if "article" in haystack:
        return "article"
    return "unknown"


def _section_heading(chunk: str) -> str:
    for line in _metadata_lines(chunk)[:4]:
        if len(line) > 100:
            continue
        if re.match(r"^(?:\d+(?:\.\d+)*\.?\s+)?[A-Z][A-Za-z0-9 ,:()/-]{3,}$", line):
            return line
    return ""


def _metadata_lines(text: str) -> list[str]:
    return [
        re.sub(r"\s+", " ", line).strip()
        for line in (text or "").splitlines()
        if re.sub(r"\s+", " ", line).strip()
    ]


def _metadata_text(value: str, limit: int = 1200) -> str:
    cleaned = re.sub(r"\s+", " ", (value or "")).strip()
    return cleaned[:limit].strip()


def _build_document_metadata(
    path: Path,
    pages: list[tuple[int | None, str]],
    pdf_pages: int,
    chunk_count: int,
    extracted_metadata: dict[str, str],
) -> dict:
    page_count = int(pdf_pages or len([page for page, text in pages if page is not None and text.strip()]))
    if page_count <= 0 and pages:
        page_count = len([text for _page, text in pages if text.strip()])
    keywords = _document_keywords(extracted_metadata, pages)
    metadata = DocumentMetadata(
        document=path.name,
        title=extracted_metadata.get("document_title") or path.stem.replace("_", " ").replace("-", " "),
        summary=_document_level_summary(path, pages, extracted_metadata, keywords),
        keywords=keywords,
        pages=max(0, page_count),
        chunks=max(0, int(chunk_count)),
        source_path=str(path),
        authors=extracted_metadata.get("authors", ""),
        document_kind=extracted_metadata.get("document_kind", "document"),
    )
    return asdict(metadata)


def _document_level_summary(
    path: Path,
    pages: list[tuple[int | None, str]],
    metadata: dict[str, str],
    keywords: list[str],
) -> str:
    abstract = _metadata_text(metadata.get("abstract", ""), limit=700)
    if abstract:
        return abstract
    sample = _metadata_text(" ".join(text for _page, text in pages[:5]), limit=3500)
    title = metadata.get("document_title") or path.stem.replace("_", " ").replace("-", " ")
    sentences = _summary_sentences(sample, title)
    if sentences:
        return _metadata_text(" ".join(sentences), limit=700)
    keyword_text = ", ".join(keywords[:8])
    if keyword_text:
        return f"Document about {keyword_text}."
    return "Indexed local document content."


def _summary_sentences(text: str, title: str) -> list[str]:
    title_key = _metadata_text(title).lower()
    candidates = re.split(r"(?<=[.!?])\s+", _metadata_text(text, limit=3500))
    sentences: list[str] = []
    for sentence in candidates:
        cleaned = _metadata_text(sentence, limit=260)
        lowered = cleaned.lower()
        if not cleaned or lowered == title_key:
            continue
        if len(cleaned) < 45 or len(cleaned.split()) < 7:
            continue
        if lowered.startswith(("abstract", "keywords", "contents", "table of contents")):
            continue
        sentences.append(cleaned)
        if len(sentences) >= 3:
            break
    return sentences


def _document_keywords(metadata: dict[str, str], pages: list[tuple[int | None, str]]) -> list[str]:
    explicit = _split_keywords(metadata.get("keywords", ""))
    if explicit:
        return explicit[:16]
    sample = " ".join(
        [
            metadata.get("document_title", ""),
            metadata.get("abstract", ""),
            " ".join(text for _page, text in pages[:3]),
        ]
    )
    terms = re.findall(r"[A-Za-z][A-Za-z0-9\-]{3,}", sample.lower())
    stopwords = {
        "about",
        "after",
        "also",
        "because",
        "been",
        "between",
        "chapter",
        "document",
        "from",
        "have",
        "into",
        "more",
        "page",
        "pages",
        "paper",
        "section",
        "that",
        "their",
        "there",
        "these",
        "this",
        "those",
        "using",
        "with",
    }
    counts: dict[str, int] = {}
    for term in terms:
        if term in stopwords or len(term) < 4:
            continue
        counts[term] = counts.get(term, 0) + 1
    return [term for term, _count in sorted(counts.items(), key=lambda item: (-item[1], item[0]))[:12]]


def _split_keywords(value: str) -> list[str]:
    return [
        _metadata_text(item, limit=80)
        for item in re.split(r"[;,|]", value or "")
        if _metadata_text(item, limit=80)
    ]


def _legacy_manifest_document_metadata(
    source_path: str,
    entry: dict,
) -> dict:
    path = Path(source_path)
    return {
        "document": path.name,
        "title": entry.get("document_title") or path.stem.replace("_", " ").replace("-", " "),
        "summary": entry.get("summary") or "",
        "keywords": entry.get("keywords") or [],
        "pages": entry.get("pdf_pages") or 0,
        "chunks": entry.get("chunks") or 0,
        "source_path": source_path,
        "authors": entry.get("authors") or "",
        "document_kind": entry.get("document_kind") or "document",
    }


def _document_metadata_from_dict(
    source_path: str,
    raw_metadata: dict,
    entry: dict,
) -> DocumentMetadata | None:
    try:
        keywords = raw_metadata.get("keywords") or []
        if isinstance(keywords, str):
            keywords = _split_keywords(keywords)
        return DocumentMetadata(
            document=str(raw_metadata.get("document") or Path(source_path).name),
            title=str(raw_metadata.get("title") or Path(source_path).stem),
            summary=str(raw_metadata.get("summary") or ""),
            keywords=[str(keyword) for keyword in keywords if str(keyword).strip()],
            pages=int(raw_metadata.get("pages") or entry.get("pdf_pages") or 0),
            chunks=int(raw_metadata.get("chunks") or entry.get("chunks") or 0),
            source_path=str(raw_metadata.get("source_path") or source_path),
            authors=str(raw_metadata.get("authors") or ""),
            document_kind=str(raw_metadata.get("document_kind") or "document"),
        )
    except Exception:
        return None


class DocumentIngestor:
    def __init__(self, settings: AppSettings) -> None:
        self.settings = settings
        self.embeddings = EmbeddingService(
            settings.embed_model,
            settings.embed_device,
            cache_dir=settings.embedding_cache_dir,
            cache_enabled=settings.embedding_cache_enabled,
        )
        self.retriever = self._create_retriever()
        self.formula_store = FormulaStore(self.settings.formula_store_path)
        self.ocr_block_store = OCRBlockStore(self.settings.ocr_block_store_path)
        self.structured_store = StructuredDocumentStore(self.settings.structured_document_store_path)

    def _create_retriever(self) -> ChromaRetriever:
        return ChromaRetriever(
            self.settings.chroma_dir,
            self.settings.collection_name,
            self.embeddings,
            settings=self.settings,
        )

    def ingest(
        self,
        reset: bool | None = None,
        progress: ProgressCallback | None = None,
    ) -> IngestResult:
        reset = self.settings.reset_db if reset is None else reset
        ensure_app_dirs(self.settings)
        with _ingest_lock(self.settings):
            backup = self._create_backup_snapshot()
            try:
                if progress and reset:
                    progress("Preparing staged rebuild", 0, 1)
                result = self._ingest_staged(progress=progress)
                if progress and reset:
                    progress("Preparing staged rebuild", 1, 1)

                rollback_reason = self._rollback_reason(result, backup)
                if rollback_reason:
                    raise IngestStateError(
                        "Ingestion rolled back because it would leave the knowledge base degraded "
                        f"({rollback_reason})."
                    )
                return result
            except Exception:
                self._restore_backup_snapshot(backup)
                raise
            finally:
                self._cleanup_backup_snapshot(backup)

    def _ingest_staged(self, progress: ProgressCallback | None) -> IngestResult:
        staging_root = Path(tempfile.mkdtemp(prefix="verilume-staging-"))
        staged_chroma = staging_root / "chroma_db"
        staged_manifest = staging_root / "ingestion_manifest.json"
        staged_formula_store = staging_root / "formulas.sqlite"
        staged_ocr_store = staging_root / "ocr_blocks.sqlite"
        staged_structured_store = staging_root / "structured_documents.sqlite"
        staged_settings = self.settings.with_overrides(
            chroma_dir=staged_chroma,
            manifest_path=staged_manifest,
            formula_store_path=staged_formula_store,
            ocr_block_store_path=staged_ocr_store,
            structured_document_store_path=staged_structured_store,
            reset_db=False,
        )
        staged = DocumentIngestor(staged_settings)
        staged.embeddings = self.embeddings
        staged.retriever = staged._create_retriever()

        try:
            result = staged._ingest_once(reset=False, progress=progress, force_all=True)
            if result.errors:
                preview = "; ".join(result.errors[:3])
                raise IngestStateError(f"Staged ingestion failed: {preview}")
            staged.retriever.close(clear_system_cache=True)
            self._install_staged_snapshot(
                staged_chroma=staged_chroma,
                staged_manifest=staged_manifest,
                staged_formula_store=staged_formula_store,
                staged_ocr_store=staged_ocr_store,
                staged_structured_store=staged_structured_store,
            )
            return result
        finally:
            shutil.rmtree(staging_root, ignore_errors=True)

    def _ingest_once(
        self,
        reset: bool,
        progress: ProgressCallback | None,
        force_all: bool,
    ) -> IngestResult:
        ensure_app_dirs(self.settings)
        if reset:
            self._reset()

        files = supported_files(self.settings.docs_dir)
        if not files:
            manifest = load_manifest(self.settings.manifest_path)
            self._remove_missing_documents(manifest)
            write_manifest(self.settings.manifest_path, manifest)
            return IngestResult(0, 0, 0, 0, 0)

        manifest = {} if force_all else load_manifest(self.settings.manifest_path)
        self._remove_missing_documents(manifest)
        manifest = {path: value for path, value in manifest.items() if Path(path).exists()}
        force_reindex = force_all or self.retriever.count() == 0
        pending: list[tuple[Path, str]] = []
        skipped = 0

        for path in files:
            digest = file_hash(path)
            key = str(path)
            entry = manifest.get(key, {})
            if not force_reindex and entry.get("hash") == digest:
                skipped += 1
                continue
            pending.append((path, digest))

        indexed_chunks = 0
        pdf_pages = 0
        errors: list[str] = []
        if progress:
            progress("Parsing documents", 0, len(pending))

        executor_class = (
            ProcessPoolExecutor
            if self.settings.process_parse_documents and len(pending) > 1
            else ThreadPoolExecutor
        )
        with executor_class(max_workers=max(1, self.settings.max_workers)) as pool:
            futures = {
                pool.submit(_parse_file, path, digest, self.settings): (path, digest)
                for path, digest in pending
            }
            for done, future in enumerate(as_completed(futures), start=1):
                path, digest = futures[future]
                try:
                    (
                        parsed_path,
                        chunks,
                        pages,
                        document_metadata,
                        formula_items,
                        ocr_blocks,
                        structured_documents,
                    ) = future.result()
                    self.retriever.delete_document(str(parsed_path))
                    self._delete_specialized_document(parsed_path.name)
                    indexed_chunks += self._index_chunks(chunks, progress)
                    self.formula_store.add_many(formula_items)
                    self.ocr_block_store.add_many(ocr_blocks)
                    for structured_document in structured_documents:
                        self.structured_store.add_structured_document(structured_document)
                    pdf_pages += pages
                    manifest[str(parsed_path)] = {
                        "hash": digest,
                        "chunks": len(chunks),
                        "pdf_pages": pages,
                        "document_metadata": document_metadata,
                    }
                except Exception as exc:
                    message = f"{path.name}: {exc}"
                    LOGGER.exception("Failed to ingest %s", path)
                    errors.append(message)
                if progress:
                    progress("Parsing documents", done, len(pending))

        if indexed_chunks:
            self.retriever.refresh_lexical_index()
        write_manifest(self.settings.manifest_path, manifest)
        return IngestResult(
            files_seen=len(files),
            files_indexed=len(pending) - len(errors),
            files_skipped=skipped,
            chunks_indexed=indexed_chunks,
            pdf_pages=pdf_pages,
            errors=errors,
        )

    def _index_chunks(
        self,
        chunks: list[DocumentChunk],
        progress: ProgressCallback | None,
    ) -> int:
        if not chunks:
            return 0

        indexed = 0
        total = len(chunks)
        start = 0
        while start < total:
            next_batch_size = _adaptive_embedding_batch_size(
                chunks[start : start + self.settings.batch_size],
                self.settings.batch_size,
            )
            batch = chunks[start : start + next_batch_size]
            texts = [chunk.text for chunk in batch]
            embeddings = self.embeddings.embed_documents(texts, batch_size=next_batch_size)
            try:
                self.retriever.add_chunks(
                    ids=[chunk.chunk_id for chunk in batch],
                    documents=texts,
                    metadatas=[_metadata(chunk) for chunk in batch],
                    embeddings=embeddings,
                )
            except Exception as exc:
                if _is_readonly_chroma_error(exc):
                    raise ReadonlyChromaError(str(exc)) from exc
                raise
            indexed += len(batch)
            if progress:
                progress("Embedding chunks", indexed, total)
            start += len(batch)
        return indexed

    def _reset(self) -> None:
        self.retriever.close(clear_system_cache=True)
        if self.settings.manifest_path.exists():
            self.settings.manifest_path.unlink()
        if self.settings.chroma_dir.exists():
            _make_tree_writable(self.settings.chroma_dir)
            shutil.rmtree(self.settings.chroma_dir, ignore_errors=True)
        self.settings.chroma_dir.mkdir(parents=True, exist_ok=True)
        _make_tree_writable(self.settings.chroma_dir)
        self.retriever = self._create_retriever()
        self.retriever.reset()
        for store_path in (
            self.settings.formula_store_path,
            self.settings.ocr_block_store_path,
            self.settings.structured_document_store_path,
        ):
            try:
                store_path.unlink(missing_ok=True)
            except OSError:
                pass
        self.formula_store = FormulaStore(self.settings.formula_store_path)
        self.ocr_block_store = OCRBlockStore(self.settings.ocr_block_store_path)
        self.structured_store = StructuredDocumentStore(self.settings.structured_document_store_path)

    def _remove_missing_documents(self, manifest: dict[str, dict]) -> None:
        missing_paths = [stored_path for stored_path in manifest if not Path(stored_path).exists()]
        for stored_path in missing_paths:
            self.retriever.delete_document(stored_path)
            self._delete_specialized_document(Path(stored_path).name)
            manifest.pop(stored_path, None)

    def _create_backup_snapshot(self) -> dict[str, object]:
        root = Path(tempfile.mkdtemp(prefix="verilume-ingest-backup-"))
        chroma_backup = root / "chroma_db"
        manifest_backup = root / "ingestion_manifest.json"
        formula_store_backup = root / "formulas.sqlite"
        ocr_store_backup = root / "ocr_blocks.sqlite"
        structured_store_backup = root / "structured_documents.sqlite"

        if self.settings.chroma_dir.exists():
            shutil.copytree(self.settings.chroma_dir, chroma_backup, dirs_exist_ok=True)
        if self.settings.manifest_path.exists():
            shutil.copy2(self.settings.manifest_path, manifest_backup)
        if self.settings.formula_store_path.exists():
            shutil.copy2(self.settings.formula_store_path, formula_store_backup)
        if self.settings.ocr_block_store_path.exists():
            shutil.copy2(self.settings.ocr_block_store_path, ocr_store_backup)
        if self.settings.structured_document_store_path.exists():
            shutil.copy2(self.settings.structured_document_store_path, structured_store_backup)

        return {
            "root": root,
            "chroma_dir": chroma_backup,
            "manifest_path": manifest_backup,
            "formula_store_path": formula_store_backup,
            "ocr_block_store_path": ocr_store_backup,
            "structured_document_store_path": structured_store_backup,
            "previous_count": self.retriever.count(),
        }

    def _restore_backup_snapshot(self, backup: dict[str, object]) -> None:
        chroma_backup = Path(str(backup["chroma_dir"]))
        manifest_backup = Path(str(backup["manifest_path"]))
        formula_store_backup = Path(str(backup["formula_store_path"]))
        ocr_store_backup = Path(str(backup["ocr_block_store_path"]))
        structured_store_backup = Path(str(backup["structured_document_store_path"]))

        self.retriever.close(clear_system_cache=True)
        if self.settings.chroma_dir.exists():
            _make_tree_writable(self.settings.chroma_dir)
            shutil.rmtree(self.settings.chroma_dir, ignore_errors=True)
        if chroma_backup.exists():
            shutil.copytree(chroma_backup, self.settings.chroma_dir, dirs_exist_ok=True)
        else:
            self.settings.chroma_dir.mkdir(parents=True, exist_ok=True)

        if manifest_backup.exists():
            shutil.copy2(manifest_backup, self.settings.manifest_path)
        elif self.settings.manifest_path.exists():
            self.settings.manifest_path.unlink()
        _restore_sqlite_file(formula_store_backup, self.settings.formula_store_path)
        _restore_sqlite_file(ocr_store_backup, self.settings.ocr_block_store_path)
        _restore_sqlite_file(structured_store_backup, self.settings.structured_document_store_path)
        self.retriever = self._create_retriever()
        self.formula_store = FormulaStore(self.settings.formula_store_path)
        self.ocr_block_store = OCRBlockStore(self.settings.ocr_block_store_path)
        self.structured_store = StructuredDocumentStore(self.settings.structured_document_store_path)

    def _install_staged_snapshot(
        self,
        *,
        staged_chroma: Path,
        staged_manifest: Path,
        staged_formula_store: Path,
        staged_ocr_store: Path,
        staged_structured_store: Path,
    ) -> None:
        self.retriever.close(clear_system_cache=True)
        if self.settings.chroma_dir.exists():
            _make_tree_writable(self.settings.chroma_dir)
            shutil.rmtree(self.settings.chroma_dir, ignore_errors=True)
        if staged_chroma.exists():
            shutil.copytree(staged_chroma, self.settings.chroma_dir, dirs_exist_ok=True)
        else:
            self.settings.chroma_dir.mkdir(parents=True, exist_ok=True)
        _make_tree_writable(self.settings.chroma_dir)

        if staged_manifest.exists():
            self.settings.manifest_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(staged_manifest, self.settings.manifest_path)
        elif self.settings.manifest_path.exists():
            self.settings.manifest_path.unlink()
        _restore_sqlite_file(staged_formula_store, self.settings.formula_store_path)
        _restore_sqlite_file(staged_ocr_store, self.settings.ocr_block_store_path)
        _restore_sqlite_file(staged_structured_store, self.settings.structured_document_store_path)
        self.retriever = self._create_retriever()
        self.formula_store = FormulaStore(self.settings.formula_store_path)
        self.ocr_block_store = OCRBlockStore(self.settings.ocr_block_store_path)
        self.structured_store = StructuredDocumentStore(self.settings.structured_document_store_path)

    def _cleanup_backup_snapshot(self, backup: dict[str, object]) -> None:
        shutil.rmtree(Path(str(backup["root"])), ignore_errors=True)

    def _rollback_reason(self, result: IngestResult, backup: dict[str, object]) -> str | None:
        if result.chunks_indexed > 0:
            self.retriever.close(clear_system_cache=True)
            self.retriever = self._create_retriever()
        current_count = self.retriever.count()
        previous_count = int(backup.get("previous_count", 0) or 0)
        if result.chunks_indexed > 0 and current_count == 0:
            return "retriever is empty after indexing"
        if result.errors and previous_count > 0 and current_count < max(1, previous_count // 2):
            return "retriever lost most indexed chunks after errors"
        return None

    def _delete_specialized_document(self, document: str) -> None:
        self.formula_store.delete_document(document)
        self.ocr_block_store.delete_document(document)
        self.structured_store.delete_document(document)


@contextmanager
def _ingest_lock(settings: AppSettings):
    lock_path = settings.chroma_dir.parent / ".verilume-ingest.lock"
    lock_path.parent.mkdir(parents=True, exist_ok=True)
    with lock_path.open("w", encoding="utf-8") as handle:
        if fcntl is not None:
            fcntl.flock(handle.fileno(), fcntl.LOCK_EX)
        try:
            yield
        finally:
            if fcntl is not None:
                fcntl.flock(handle.fileno(), fcntl.LOCK_UN)


def _restore_sqlite_file(source: Path, target: Path) -> None:
    target.parent.mkdir(parents=True, exist_ok=True)
    if source.exists():
        shutil.copy2(source, target)
        return
    try:
        target.unlink(missing_ok=True)
    except OSError:
        pass


def _metadata(chunk: DocumentChunk) -> dict:
    metadata = dict(chunk.metadata)
    metadata.update(
        {
            "document": chunk.document,
            "source_path": str(chunk.source_path),
            "page": chunk.page or 0,
            "chunk_index": chunk.chunk_index,
            "file_hash": chunk.file_hash,
        }
    )
    return metadata


def _adaptive_embedding_batch_size(chunks: list[DocumentChunk], base_batch_size: int) -> int:
    base = max(1, int(base_batch_size))
    if not chunks:
        return base
    average_chars = sum(len(chunk.text or "") for chunk in chunks) / max(1, len(chunks))
    if average_chars <= 450:
        return min(512, base * 2)
    if average_chars >= 2200:
        base = max(16, base // 2)

    available_memory = _available_memory_bytes()
    if available_memory is None:
        return base
    estimated_bytes_per_chunk = max(512, int(average_chars * 4)) + 1536 * 4
    memory_safe_batch = max(1, int(available_memory * 0.15) // estimated_bytes_per_chunk)
    return max(1, min(base, memory_safe_batch))


def _available_memory_bytes() -> int | None:
    if hasattr(os, "sysconf"):
        try:
            pages = int(os.sysconf("SC_AVPHYS_PAGES"))
            page_size = int(os.sysconf("SC_PAGE_SIZE"))
            return pages * page_size
        except (OSError, ValueError):
            return None
    return None


def _is_readonly_chroma_error(exc: Exception) -> bool:
    text = str(exc).lower()
    return "readonly database" in text or "read-only database" in text


def _make_tree_writable(path: Path) -> None:
    if not path.exists():
        return
    for item in [path, *path.rglob("*")]:
        try:
            mode = item.stat().st_mode
            item.chmod(mode | 0o200)
        except OSError:
            pass
