from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from verilume.core.schemas import DocumentChunk
from verilume.ingest import (
    DocumentIngestor,
    _adaptive_embedding_batch_size,
    _build_document_metadata,
    _extract_document_metadata,
    _normalize_pdf_text,
    chunk_text_semantic,
    document_metadata_from_manifest,
    extract_pages,
    load_manifest,
    removable_documents,
    remove_documents,
    supported_extensions,
    write_manifest,
)
from verilume.settings import AppSettings


class FakeRetriever:
    def __init__(self) -> None:
        self.deleted_paths: list[str] = []

    def delete_document(self, source_path: str) -> None:
        self.deleted_paths.append(source_path)


class FakeRemovalIngestor:
    last_retriever: FakeRetriever | None = None

    def __init__(self, settings: AppSettings) -> None:
        self.retriever = FakeRetriever()
        FakeRemovalIngestor.last_retriever = self.retriever


def _chunk(text: str) -> DocumentChunk:
    return DocumentChunk(
        chunk_id=f"chunk-{len(text)}",
        text=text,
        source_path=Path("doc.pdf"),
        document="doc.pdf",
        page=1,
        chunk_index=0,
        file_hash="hash",
    )


class IngestCleanupTests(unittest.TestCase):
    def test_supported_extensions_include_presentations_and_images(self) -> None:
        extensions = supported_extensions()

        self.assertTrue(
            {".pptx", ".png", ".jpg", ".jpeg", ".tiff", ".webp"}.issubset(extensions)
        )

    def test_document_metadata_is_extracted_from_research_text(self) -> None:
        text = """
        Replica Exchange Hamiltonian Monte Carlo for Hydrological Models
        Alex Jordan Sample
        University of Luxembourg

        Abstract
        Replica Exchange Hamiltonian Monte Carlo combines Hamiltonian Monte Carlo with
        replica exchange to improve exploration of multimodal posterior distributions
        in Bayesian hydrological modelling.

        Keywords: Hamiltonian Monte Carlo; Replica Exchange; Bayesian inference

        1 Introduction
        The method is evaluated on HBV models.
        """

        metadata = _extract_document_metadata(Path("hremc-paper.pdf"), [(1, text)])

        self.assertEqual(
            metadata["document_title"],
            "Replica Exchange Hamiltonian Monte Carlo for Hydrological Models",
        )
        self.assertIn("Alex Jordan Sample", metadata["authors"])
        self.assertIn("combines Hamiltonian Monte Carlo", metadata["abstract"])
        self.assertIn("Replica Exchange", metadata["keywords"])
        self.assertEqual(metadata["document_kind"], "scientific_paper")

    def test_document_level_summary_is_manifest_ready(self) -> None:
        text = """
        Bayesian Statistics (2012)

        Abstract
        Comprehensive Bayesian statistics textbook covering regression, Gibbs sampling,
        Markov chain Monte Carlo, model assessment, diagnostics, and importance sampling.

        Keywords: Bayesian inference; Gibbs sampling; MCMC; regression
        """
        path = Path("bayesian-statistics.pdf")
        extracted = _extract_document_metadata(path, [(1, text)])

        metadata = _build_document_metadata(
            path,
            [(1, text)],
            pdf_pages=350,
            chunk_count=120,
            extracted_metadata=extracted,
        )

        self.assertEqual(metadata["title"], "Bayesian Statistics (2012)")
        self.assertIn("Gibbs sampling", metadata["summary"])
        self.assertIn("MCMC", metadata["keywords"])
        self.assertEqual(metadata["pages"], 350)
        self.assertEqual(metadata["chunks"], 120)

    def test_document_metadata_from_manifest_reads_document_summaries(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            root = Path(tmp_dir)
            docs_dir = root / "docs"
            docs_dir.mkdir()
            manifest_path = root / "manifest.json"
            document_path = docs_dir / "dic.pdf"
            write_manifest(
                manifest_path,
                {
                    str(document_path): {
                        "hash": "abc",
                        "chunks": 8,
                        "pdf_pages": 12,
                        "document_metadata": {
                            "document": "dic.pdf",
                            "title": "Dictionary",
                            "summary": "Terminology and definitions.",
                            "keywords": ["terminology", "definitions"],
                            "pages": 12,
                            "chunks": 8,
                            "source_path": str(document_path),
                        },
                    }
                },
            )
            settings = AppSettings(docs_dir=docs_dir, manifest_path=manifest_path)

            metadata = document_metadata_from_manifest(settings)

        self.assertEqual(len(metadata), 1)
        self.assertEqual(metadata[0].document, "dic.pdf")
        self.assertEqual(metadata[0].summary, "Terminology and definitions.")

    def test_adaptive_embedding_batch_size_uses_chunk_length(self) -> None:
        self.assertEqual(
            _adaptive_embedding_batch_size([_chunk("short text")] * 4, 128),
            256,
        )
        self.assertEqual(
            _adaptive_embedding_batch_size([_chunk("medium text " * 90)] * 4, 128),
            128,
        )
        self.assertEqual(
            _adaptive_embedding_batch_size([_chunk("long text " * 260)] * 4, 128),
            64,
        )

    def test_semantic_chunking_prefers_sentence_boundaries(self) -> None:
        text = (
            "Alpha sentence one. Beta sentence two carries important context. "
            "Gamma sentence three closes the thought."
        )

        chunks = chunk_text_semantic(text, chunk_size=62, chunk_overlap=0)

        self.assertGreaterEqual(len(chunks), 2)
        self.assertTrue(all(chunk.endswith(".") for chunk in chunks))
        self.assertIn("Beta sentence two", " ".join(chunks))

    def test_normalize_pdf_text_removes_icon_font_fragments(self) -> None:
        text = (
            "Alex Jordan Sample\n"
            "Rue Example, L-1234 Luxembourg | /ne+352 600 000 000 | "
            "alex.sample@example.com | /gtb| /♀nedn\n"
            "quanti-\n"
            "tative methods"
        )

        normalized = _normalize_pdf_text(text)

        self.assertIn("Rue Example, L-1234 Luxembourg | +352 600 000 000 | alex.sample@example.com", normalized)
        self.assertNotIn("/ne", normalized)
        self.assertNotIn("/gtb", normalized)
        self.assertNotIn("/♀nedn", normalized)
        self.assertIn("quantitative methods", normalized)

    def test_missing_documents_are_removed_from_manifest_and_vector_store(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            existing_path = tmp_path / "exists.pdf"
            missing_path = tmp_path / "missing.pdf"
            existing_path.write_text("hello", encoding="utf-8")

            settings = AppSettings(
                docs_dir=tmp_path,
                chroma_dir=tmp_path / "chroma",
                manifest_path=tmp_path / "manifest.json",
            )
            ingestor = DocumentIngestor(settings)
            ingestor.retriever = FakeRetriever()

            manifest = {
                str(existing_path): {"hash": "a"},
                str(missing_path): {"hash": "b"},
            }
            ingestor._remove_missing_documents(manifest)

            self.assertEqual(list(manifest), [str(existing_path)])
            self.assertEqual(ingestor.retriever.deleted_paths, [str(missing_path)])

    def test_removable_documents_and_remove_documents_stay_in_sync(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            docs_dir = tmp_path / "docs"
            nested_dir = docs_dir / "profiles"
            nested_dir.mkdir(parents=True)
            keep_path = docs_dir / "keep.pdf"
            remove_path = nested_dir / "remove.pdf"
            keep_path.write_text("keep", encoding="utf-8")
            remove_path.write_text("remove", encoding="utf-8")

            settings = AppSettings(
                docs_dir=docs_dir,
                chroma_dir=tmp_path / "chroma",
                manifest_path=tmp_path / "manifest.json",
            )
            write_manifest(
                settings.manifest_path,
                {
                    str(keep_path): {"hash": "keep", "chunks": 2},
                    str(remove_path): {"hash": "remove", "chunks": 1},
                },
            )

            self.assertEqual(removable_documents(docs_dir), ["keep.pdf", "profiles/remove.pdf"])

            with patch("verilume.ingest.DocumentIngestor", new=FakeRemovalIngestor):
                removed = remove_documents(settings, ["profiles/remove.pdf"])

            self.assertEqual(removed, ["profiles/remove.pdf"])
            self.assertFalse(remove_path.exists())
            self.assertTrue(keep_path.exists())
            self.assertEqual(list(load_manifest(settings.manifest_path)), [str(keep_path)])
            expected_deleted_paths = [str(remove_path)]
            resolved_remove_path = str(remove_path.resolve())
            if resolved_remove_path != str(remove_path):
                expected_deleted_paths.append(resolved_remove_path)
            self.assertEqual(
                FakeRemovalIngestor.last_retriever.deleted_paths,
                expected_deleted_paths,
            )

    def test_image_files_are_ocrd_through_image_handler(self) -> None:
        from PIL import Image

        with tempfile.TemporaryDirectory() as tmp_dir:
            image_path = Path(tmp_dir) / "scan.png"
            Image.new("RGB", (240, 80), "white").save(image_path)

            with patch("verilume.ingest._ocr_pil_image", return_value="Scanned invoice total"):
                pages, pdf_pages = extract_pages(image_path)

            self.assertEqual(pages, [(1, "Scanned invoice total")])
            self.assertEqual(pdf_pages, 0)

    def test_scanned_pdf_pages_fall_back_to_ocr(self) -> None:
        from PIL import Image

        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            image_path = tmp_path / "scan.png"
            pdf_path = tmp_path / "scan.pdf"

            image = Image.new("RGB", (240, 80), "white")
            image.save(image_path)

            pdf = canvas.Canvas(str(pdf_path), pagesize=(240, 80))
            pdf.drawImage(ImageReader(str(image_path)), 0, 0, width=240, height=80)
            pdf.showPage()
            pdf.save()

            with patch("verilume.ingest._ocr_pdf_page", return_value="Scanned PDF OCR text"):
                pages, pdf_pages = extract_pages(pdf_path)

            self.assertEqual(pages, [(1, "Scanned PDF OCR text")])
            self.assertEqual(pdf_pages, 1)

    def test_powerpoint_slides_are_extracted(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp_dir:
            pptx_path = Path(tmp_dir) / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Quarterly Research Update"
            slide.placeholders[1].text = "Pilot results improved document recall by 12 percent."
            presentation.save(pptx_path)

            pages, pdf_pages = extract_pages(pptx_path)

            self.assertEqual(pdf_pages, 0)
            self.assertEqual(len(pages), 1)
            self.assertEqual(pages[0][0], 1)
            self.assertIn("Quarterly Research Update", pages[0][1])
            self.assertIn("document recall by 12 percent", pages[0][1])
